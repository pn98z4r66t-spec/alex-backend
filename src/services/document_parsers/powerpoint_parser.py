"""
PowerPoint Presentation Parser
Extracts text, notes, and metadata from PowerPoint files (.pptx)
"""
from pptx import Presentation
from typing import Dict, List, Any
import logging

logger = logging.getLogger(__name__)


class PowerPointParser:
    """Parser for PowerPoint presentations (.pptx)"""
    
    @staticmethod
    def parse(file_path: str) -> Dict[str, Any]:
        """
        Parse PowerPoint file and extract content
        
        Args:
            file_path: Path to .pptx file
            
        Returns:
            Dictionary containing:
                - text: Full text content
                - slides: List of slide contents
                - metadata: Presentation properties
                - has_images: Whether presentation contains images
                - has_tables: Whether presentation contains tables
        """
        try:
            result = {
                'text': '',
                'slides': [],
                'metadata': {},
                'has_images': False,
                'has_tables': False,
                'success': True,
                'error': None
            }
            
            prs = Presentation(file_path)
            
            # Extract metadata
            if hasattr(prs, 'core_properties'):
                props = prs.core_properties
                result['metadata'] = {
                    'title': props.title or '',
                    'author': props.author or '',
                    'subject': props.subject or '',
                    'keywords': props.keywords or '',
                    'created': str(props.created) if props.created else '',
                    'modified': str(props.modified) if props.modified else '',
                    'last_modified_by': props.last_modified_by or '',
                }
            
            # Get slide dimensions
            result['metadata']['slide_width'] = prs.slide_width
            result['metadata']['slide_height'] = prs.slide_height
            result['metadata']['slide_count'] = len(prs.slides)
            
            # Process each slide
            all_text = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_data = {
                    'slide_number': slide_idx,
                    'title': '',
                    'content': [],
                    'notes': '',
                    'has_images': False,
                    'has_tables': False,
                    'shape_count': len(slide.shapes)
                }
                
                # Extract content from shapes
                for shape in slide.shapes:
                    # Text content
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        
                        # Try to identify title
                        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                            if shape.placeholder_format.type == 1:  # Title placeholder
                                slide_data['title'] = text
                            else:
                                slide_data['content'].append(text)
                        else:
                            slide_data['content'].append(text)
                    
                    # Check for tables
                    if hasattr(shape, 'table'):
                        slide_data['has_tables'] = True
                        result['has_tables'] = True
                        
                        # Extract table data
                        table = shape.table
                        table_text = []
                        for row in table.rows:
                            row_text = ' | '.join([cell.text.strip() for cell in row.cells])
                            table_text.append(row_text)
                        
                        slide_data['content'].append(f"[Table]\n" + '\n'.join(table_text))
                    
                    # Check for images
                    if hasattr(shape, 'image'):
                        slide_data['has_images'] = True
                        result['has_images'] = True
                
                # Extract speaker notes
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text.strip():
                        slide_data['notes'] = notes_frame.text.strip()
                
                result['slides'].append(slide_data)
                
                # Add slide content to text
                slide_text = f"\n[Slide {slide_idx}]"
                if slide_data['title']:
                    slide_text += f"\nTitle: {slide_data['title']}"
                if slide_data['content']:
                    slide_text += f"\nContent:\n" + '\n'.join(slide_data['content'])
                if slide_data['notes']:
                    slide_text += f"\nNotes: {slide_data['notes']}"
                
                all_text.append(slide_text)
            
            # Combine all text
            result['text'] = '\n\n'.join(all_text)
            
            # Add statistics
            result['statistics'] = {
                'total_slides': len(result['slides']),
                'slides_with_titles': sum(1 for s in result['slides'] if s['title']),
                'slides_with_notes': sum(1 for s in result['slides'] if s['notes']),
                'slides_with_images': sum(1 for s in result['slides'] if s['has_images']),
                'slides_with_tables': sum(1 for s in result['slides'] if s['has_tables']),
                'total_characters': len(result['text']),
                'total_words': len(result['text'].split()),
            }
            
            return result
            
        except Exception as e:
            logger.error(f"Error parsing PowerPoint file: {e}")
            return {
                'success': False,
                'error': str(e),
                'text': '',
                'slides': [],
                'metadata': {}
            }
    
    @staticmethod
    def get_summary(parsed_data: Dict[str, Any]) -> str:
        """
        Generate a human-readable summary of the PowerPoint presentation
        
        Args:
            parsed_data: Output from parse() method
            
        Returns:
            Summary string
        """
        if not parsed_data.get('success'):
            return f"Failed to parse PowerPoint file: {parsed_data.get('error')}"
        
        summary_parts = []
        
        # Basic info
        stats = parsed_data.get('statistics', {})
        summary_parts.append(
            f"PowerPoint Presentation with {stats.get('total_slides', 0)} slides"
        )
        
        # Metadata
        metadata = parsed_data.get('metadata', {})
        if metadata.get('title'):
            summary_parts.append(f"Title: {metadata['title']}")
        if metadata.get('author'):
            summary_parts.append(f"Author: {metadata['author']}")
        
        # Content statistics
        summary_parts.append(
            f"Content: {stats.get('total_words', 0)} words, "
            f"{stats.get('slides_with_titles', 0)} titled slides"
        )
        
        # Additional features
        features = []
        if stats.get('slides_with_notes', 0) > 0:
            features.append(f"{stats['slides_with_notes']} slides with notes")
        if stats.get('slides_with_images', 0) > 0:
            features.append(f"{stats['slides_with_images']} slides with images")
        if stats.get('slides_with_tables', 0) > 0:
            features.append(f"{stats['slides_with_tables']} slides with tables")
        
        if features:
            summary_parts.append("Features: " + ', '.join(features))
        
        return '\n'.join(summary_parts)

